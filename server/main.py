"""
FastAPI RAG Server for Sales Coach Extension
Uses Ollama for embeddings and PostgreSQL for vector storage
"""
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional
import ollama
import database
import io

app = FastAPI(title="Sales Coach RAG API", version="1.0.0")

# Enable CORS for Chrome extension
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Ollama embedding model
EMBED_MODEL = "nomic-embed-text"
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100

# Pydantic models
class DocumentInput(BaseModel):
    title: str
    content: str

class SearchQuery(BaseModel):
    query: str
    limit: Optional[int] = 3

class DocumentResponse(BaseModel):
    id: int
    title: str
    content: str
    similarity: Optional[float] = None

# Utility functions
def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks"""
    chunks = []
    start = 0
    text = text.strip()
    
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        
        # Try to break at sentence or word boundary
        if end < len(text):
            # Look for sentence end
            last_period = chunk.rfind('.')
            last_newline = chunk.rfind('\n')
            break_point = max(last_period, last_newline)
            
            if break_point > chunk_size // 2:
                chunk = text[start:start + break_point + 1]
                end = start + break_point + 1
        
        chunks.append(chunk.strip())
        start = end - overlap
    
    return [c for c in chunks if c]  # Remove empty chunks

def get_embedding(text: str) -> List[float]:
    """Generate embedding using Ollama nomic-embed-text"""
    try:
        response = ollama.embeddings(model=EMBED_MODEL, prompt=text)
        return response['embedding']
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Embedding error: {str(e)}")

# API Endpoints
@app.on_event("startup")
async def startup():
    """Initialize database on startup"""
    try:
        database.init_database()
    except Exception as e:
        print(f"⚠️ Database init warning: {e}")

@app.get("/")
async def root():
    """Health check"""
    return {"status": "ok", "service": "Sales Coach RAG API"}

@app.get("/health")
async def health():
    """Health check with stats"""
    try:
        count = database.get_document_count()
        return {"status": "healthy", "documents": count}
    except Exception as e:
        return {"status": "error", "error": str(e)}

@app.post("/documents")
async def add_document(doc: DocumentInput):
    """Add a document to the knowledge base"""
    try:
        # Chunk the document
        chunks = chunk_text(doc.content)
        
        if not chunks:
            raise HTTPException(status_code=400, detail="Document is empty or too short")
        
        # Process each chunk
        doc_ids = []
        for i, chunk in enumerate(chunks):
            embedding = get_embedding(chunk)
            doc_id = database.add_document(doc.title, chunk, i, embedding)
            doc_ids.append(doc_id)
        
        return {
            "success": True,
            "title": doc.title,
            "chunks": len(chunks),
            "ids": doc_ids
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/documents")
async def list_documents():
    """List all documents in knowledge base"""
    try:
        docs = database.get_all_documents()
        return {"documents": docs, "count": len(docs)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/documents/{title}")
async def delete_document(title: str):
    """Delete a document by title"""
    try:
        deleted = database.delete_document(title)
        if deleted == 0:
            raise HTTPException(status_code=404, detail="Document not found")
        return {"success": True, "deleted_chunks": deleted}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/search")
async def search_documents(query: SearchQuery):
    """Search for relevant documents"""
    try:
        # Generate query embedding
        query_embedding = get_embedding(query.query)
        
        # Search similar documents
        results = database.search_similar(query_embedding, query.limit)
        
        return {
            "query": query.query,
            "results": results,
            "count": len(results)
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/embed")
async def generate_embedding(text: str):
    """Generate embedding for text"""
    try:
        embedding = get_embedding(text)
        return {"embedding": embedding, "dimensions": len(embedding)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/extract-file")
async def extract_file(file: UploadFile = File(...)):
    """Extract text content from PDF, DOCX, or PPTX files"""
    try:
        filename = file.filename.lower()
        content = await file.read()
        
        if filename.endswith('.pdf'):
            return await extract_pdf_content(content)
        elif filename.endswith('.docx'):
            return await extract_docx_content(content)
        elif filename.endswith('.pptx'):
            return await extract_pptx_content(content)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type. Use PDF, DOCX, or PPTX")
            
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Extraction error: {str(e)}")

async def extract_pdf_content(content: bytes):
    """Extract text from PDF"""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        
        text_content = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
        
        extracted_text = "\n\n".join(text_content)
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from PDF")
        
        return {"content": extracted_text, "pages": len(pdf_reader.pages), "type": "pdf"}
        
    except ImportError:
        raise HTTPException(status_code=500, detail="PyPDF2 not installed")

async def extract_docx_content(content: bytes):
    """Extract text from DOCX (Word document)"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        extracted_text = "\n\n".join(text_content)
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from DOCX")
        
        return {"content": extracted_text, "paragraphs": len(doc.paragraphs), "type": "docx"}
        
    except ImportError:
        raise HTTPException(status_code=500, detail="python-docx not installed")

async def extract_pptx_content(content: bytes):
    """Extract text from PPTX (PowerPoint)"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        
        text_content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # More than just the header
                text_content.append("\n".join(slide_text))
        
        extracted_text = "\n\n".join(text_content)
        
        if not extracted_text.strip():
            raise HTTPException(status_code=400, detail="Could not extract text from PPTX")
        
        return {"content": extracted_text, "slides": len(prs.slides), "type": "pptx"}
        
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed")

# Keep old endpoint for backwards compatibility
@app.post("/extract-pdf")
async def extract_pdf_legacy(file: UploadFile = File(...)):
    """Legacy PDF extraction endpoint"""
    content = await file.read()
    return await extract_pdf_content(content)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
